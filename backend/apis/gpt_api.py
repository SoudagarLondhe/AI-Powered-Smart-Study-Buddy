# backend/apis/gpt_api.py
import os
import json
import pathlib
from tempfile import NamedTemporaryFile
from typing import Callable, List, Optional

from dotenv import load_dotenv
from fastapi import HTTPException, UploadFile, File, Form, Body
from sqlalchemy import select, func, delete
from sqlalchemy.orm import Session

# File parsers
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

# DB models
from models import Course, Summary
from fastapi import Query

# OpenAI for AI-powered summarization
from openai import OpenAI

load_dotenv()

# --- Global session factory ---
_SESSION_FACTORY: Optional[Callable[[], Session]] = None
def set_session_factory(factory: Callable[[], Session]) -> None:
    global _SESSION_FACTORY
    _SESSION_FACTORY = factory

# --- Helper responses ---
def _fail(msg: str) -> dict:
    return {"status": "FAIL", "statusCode": 200, "message": msg, "data": ""}

def _success(data_str: str, message: str = "") -> dict:
    return {"status": "SUCCESS", "statusCode": 200, "message": message, "data": data_str}

# --- File extraction helpers ---
MAX_BYTES = 50 * 1024 * 1024  # 50 MB

def _read_pdf(path: str) -> str:
    parts: List[str] = []
    with pdfplumber.open(path) as pdf:
        for p in pdf.pages:
            parts.append(p.extract_text() or "")
    return "\n".join(parts).strip()

def _read_docx(path: str) -> str:
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs).strip()

def _read_pptx(path: str) -> str:
    prs = Presentation(path)
    parts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts).strip()

def _extract_text_from_upload(upload: UploadFile) -> str:
    name = upload.filename or "uploaded"
    _, ext = os.path.splitext(name.lower())

    upload.file.seek(0, os.SEEK_END)
    size = upload.file.tell()
    upload.file.seek(0)
    if size > MAX_BYTES:
        raise HTTPException(status_code=413, detail=f"{name} exceeds {MAX_BYTES // (1024*1024)}MB limit")

    with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(upload.file.read())
        tmp_path = tmp.name

    try:
        if ext == ".pdf":
            text = _read_pdf(tmp_path)
        elif ext == ".docx":
            text = _read_docx(tmp_path)
        elif ext == ".pptx":
            text = _read_pptx(tmp_path)
        else:
            raise HTTPException(status_code=415, detail=f"Unsupported file type '{ext}'. Use PDF, DOCX, or PPTX.")
    finally:
        try:
            os.remove(tmp_path)
        except Exception:
            pass

    if not text:
        raise HTTPException(status_code=400, detail=f"No readable text found in {name}. If it is a scanned PDF, add OCR.")
    return text

def _extract_many(uploads: List[UploadFile]) -> List[tuple[str, str]]:
    out: List[tuple[str, str]] = []
    for up in uploads:
        text = _extract_text_from_upload(up)
        stem = pathlib.Path(up.filename or "uploaded").stem
        out.append((stem, text))
    return out

def _derive_course_name(stems: List[str], override: Optional[str]) -> str:
    if override:
        return override[:255]
    if not stems:
        return "Untitled"
    base = " & ".join(stems[:3])
    if len(stems) > 3:
        base += f" (+{len(stems)-3} more)"
    return base[:255]

# --- POST /addcourse ---
def ingest_and_store_endpoint(
    files: Optional[List[UploadFile]] = File(default=None),
    course_name: Optional[str] = Form(default=None),
    body: Optional[dict] = Body(default=None),
):
    if _SESSION_FACTORY is None:
        return _fail("Server misconfigured: no DB session factory is set.")

    try:
        if files is not None and len(files) > 0:
            extracted = _extract_many(files)
            stems = [s for s, _ in extracted]
            combined_content = "\n\n".join(
                [f"=== FILE: {s} ===\n{t}" for s, t in extracted]
            ).strip()
            if not combined_content:
                return _fail("No readable text found in uploaded files.")

            final_name = _derive_course_name(stems, course_name)
            with _SESSION_FACTORY() as db:
                row = Course(course_name=final_name, course_content=combined_content)
                db.add(row)
                db.flush()
                new_id = row.course_id
                db.commit()

            return _success(
                f"Saved 1 course(s): {final_name}=>id={new_id}",
                "Stored concatenated text from uploaded files into a single course."
            )

        elif body and isinstance(body, dict) and body.get("content"):
            cname = (body.get("course_name") or "Untitled")[:255]
            content = str(body["content"]).strip()
            if not content:
                return _fail("Content is empty.")

            with _SESSION_FACTORY() as db:
                row = Course(course_name=cname, course_content=content)
                db.add(row)
                db.flush()
                new_id = row.course_id
                db.commit()

            return _success(
                f"Saved 1 course(s): {cname}=>id={new_id}",
                "Stored raw text into courses table."
            )

        else:
            return _fail("Provide 'files' (one or many) OR JSON {content, course_name}.")

    except HTTPException:
        raise
    except Exception as e:
        return _fail(f"Ingestion failed: {type(e).__name__}")

# --- GET /courses ---
def list_courses(limit: int = 25, offset: int = 0, q: Optional[str] = None):
    if _SESSION_FACTORY is None:
        return _fail("Server misconfigured: no DB session factory is set.")

    with _SESSION_FACTORY() as db:
        stmt = select(
            Course.course_id,
            Course.course_name,
            func.length(Course.course_content).label("content_len"),
        )
        if q:
            from sqlalchemy import func as _func
            stmt = stmt.where(_func.lower(Course.course_name).like(f"%{q.lower()}%"))
        stmt = stmt.order_by(Course.course_id.desc()).offset(offset).limit(limit)

        rows = db.execute(stmt).all()
        payload = [
            {"course_id": r.course_id, "course_name": r.course_name, "content_len": int(r.content_len or 0)}
            for r in rows
        ]
        return _success(json.dumps(payload, ensure_ascii=False), message=f"Fetched {len(payload)} course(s).")

# --- GET /courses/{course_id} ---
def get_course(course_id: int):
    if _SESSION_FACTORY is None:
        return _fail("Server misconfigured: no DB session factory is set.")

    with _SESSION_FACTORY() as db:
        row = db.execute(
            select(Course.course_id, Course.course_name, Course.course_content)
            .where(Course.course_id == course_id)
        ).one_or_none()

        if not row:
            return _fail(f"Course id={course_id} not found")

        course_id_db, course_name_db, content = row
        payload = {
            "course_id": course_id_db,
            "course_name": course_name_db,
            "course_content": content,
        }
        return _success(json.dumps(payload, ensure_ascii=False), message=f"Fetched course id={course_id}.")

# --- POST /courses/{course_id}/summary ---
def generate_course_summary(course_id: int, body: dict = Body(...)):
    """
    Generate or replace a summary for a given course_id and length: short|medium|long
    """
    if _SESSION_FACTORY is None:
        return _fail("Server misconfigured: no DB session factory is set.")

    summary_length = body.get("summary_length")
    summary_map = {"short": 300, "medium": 700, "long": 1200}
    max_chars = summary_map.get(summary_length.lower())
    if not max_chars:
        return _fail("Invalid summary_length. Use one of: short, medium, long.")

    with _SESSION_FACTORY() as db:
        row = db.execute(
            select(Course.course_id, Course.course_name, Course.course_content)
            .where(Course.course_id == course_id)
        ).one_or_none()
        if not row:
            return _fail(f"Course id={course_id} not found")

        _, cname, content = row
        content = (content or "").strip()
        if not content:
            return _fail(f"Course id={course_id} has no content")

    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    system_prompt = (
        f"You are a helpful teaching assistant. Summarize the following course material in {summary_length} form. "
        "Explain concepts in very simple, clear language that any student can understand. "
        "Avoid jargon. Use bullet points and short sentences."
    )

    user_prompt = f"Summarize this course content:\n\n{content[:4000]}"

    try:
        resp = client.responses.create(
            model="gpt-4o-mini",
            input=f"{system_prompt}\n\n{user_prompt}",
            max_output_tokens=max_chars,
        )
        summary_text = getattr(resp, "output_text", None) or str(resp)
    except Exception as e:
        return _fail(f"AI summarization failed: {type(e).__name__}")

    with _SESSION_FACTORY() as db:
        db.execute(
            delete(Summary).where(Summary.course_id == course_id, Summary.summary_length == summary_length)
        )
        db.add(Summary(course_id=course_id, summary_length=summary_length, summary_content=summary_text))
        db.commit()

    return _success(
        summary_text,
        message=f"Summary ({summary_length}) generated and stored for course_id={course_id}."
    )


def get_course_summary(
    course_id: int,
    summary_length: str = Query(..., description="short | medium | long"),
):
    """
    GET /courses/{course_id}/summary?summary_length=short|medium|long
    Returns the stored summary for (course_id, summary_length).
    """
    if _SESSION_FACTORY is None:
        return _fail("Server misconfigured: no DB session factory is set.")

    length = (summary_length or "").strip().lower()
    if length not in {"short", "medium", "long"}:
        return _fail("Invalid summary_length. Use one of: short, medium, long.")

    with _SESSION_FACTORY() as db:
        row = db.execute(
            select(Summary.summary_id, Summary.summary_content)
            .where(Summary.course_id == course_id, Summary.summary_length == length)
        ).one_or_none()

        if not row:
            return _fail(
                f"No summary found for course_id={course_id} and length='{length}'. "
                f"POST /courses/{course_id}/summary to generate one."
            )

        summary_id, summary_content = row
        payload = {
            "summary_id": summary_id,
            "course_id": course_id,
            "summary_length": length,
            "summary_content": summary_content,
        }
        return _success(
            data_str=json.dumps(payload, ensure_ascii=False),
            message=f"Fetched summary for course_id={course_id}, length={length}."
        )